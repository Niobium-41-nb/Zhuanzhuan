from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

GREEN = RGBColor(0x2D, 0x6A, 0x4F)
DARK = RGBColor(0x1B, 0x43, 0x32)
ACCENT = RGBColor(0xE7, 0x6F, 0x51)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xFD, 0xFB, 0xF7)
TEXT = RGBColor(0x1F, 0x29, 0x37)
MUTED = RGBColor(0x6B, 0x72, 0x80)
W = prs.slide_width
H = prs.slide_height
LOGO = "F:/zhuanzhuan/Zhuanzhuan/docs/工作日报/hzau-logo.png"

def logo(s):
    s.shapes.add_picture(LOGO, W-Inches(1.6), Inches(0.15), width=Inches(1.4), height=Inches(0.4))

def bar(s):
    r = s.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(0.04))
    r.fill.solid(); r.fill.fore_color.rgb = GREEN; r.line.fill.background()
    r = s.shapes.add_shape(1, Inches(0), H-Inches(0.03), W, Inches(0.03))
    r.fill.solid(); r.fill.fore_color.rgb = GREEN; r.line.fill.background()

def pn(s, n):
    t = s.shapes.add_textbox(W-Inches(1.2), H-Inches(0.45), Inches(1), Inches(0.3))
    p = t.text_frame.paragraphs[0]
    p.text = f"{n}/13"; p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(0x9C, 0xA3, 0xAF); p.alignment = PP_ALIGN.RIGHT

def title(s, txt):
    t = s.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11), Inches(0.5))
    p = t.text_frame.paragraphs[0]
    p.text = txt; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = TEXT

def sec(s, x, y, w, title_txt, items, color=GREEN):
    t = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.4))
    p = t.text_frame.paragraphs[0]
    p.text = title_txt; p.font.size = Pt(17); p.font.bold = True; p.font.color.rgb = color
    t = s.shapes.add_textbox(Inches(x), Inches(y+0.45), Inches(w), Inches(len(items)*0.3+0.3))
    tf = t.text_frame; first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = f"  {item}"; p.font.size = Pt(13); p.font.color.rgb = MUTED; p.space_after = Pt(4)

# === Page 1: Cover ===
s = prs.slides.add_slide(prs.slide_layouts[6])
r = s.shapes.add_shape(1, Inches(0), Inches(0), W, H)
r.fill.solid(); r.fill.fore_color.rgb = DARK; r.line.fill.background()
logo(s)
t = s.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.33), Inches(3))
tf = t.text_frame
p = tf.paragraphs[0]; p.text = "转转"; p.font.size = Pt(60); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph(); p.text = "校园二手物品交易平台"; p.font.size = Pt(24); p.font.color.rgb = RGBColor(0xE0,0xE0,0xE0); p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph(); p.text = ""; p.font.size = Pt(14)
p = tf.add_paragraph(); p.text = "软件工程实训 · 课程设计答辩"; p.font.size = Pt(18); p.font.color.rgb = RGBColor(0xBB,0xBB,0xBB); p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph(); p.text = ""; p.font.size = Pt(12)
p = tf.add_paragraph(); p.text = "张可凡 · 阮崇轩 · 李硕 · 刘轩霆 · 胡欣悦"; p.font.size = Pt(15); p.font.color.rgb = RGBColor(0x99,0x99,0x99); p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph(); p.text = "华中农业大学 · 2026年7月"; p.font.size = Pt(13); p.font.color.rgb = RGBColor(0x88,0x88,0x88); p.alignment = PP_ALIGN.CENTER

# === Page 2: Demo - Browse & Search ===
s = prs.slides.add_slide(prs.slide_layouts[6])
bar(s); logo(s); pn(s, 2); title(s, "项目演示 — 浏览与搜索")
sec(s, 0.8, 1.1, 5.5, "首页展示", [
    "Hero 横幅 + 品牌标语 + 实时统计数据",
    "搜索框快速跳转商品列表",
    "公告栏水平滚动轮播 (CSS Animation)",
    "8 个分类 Tab (Emoji) 切换加载",
    "商品网格响应式布局 + 骨架屏加载",
    "卡片 hover 放大 + 阴影 + 入场动画"])
sec(s, 7, 1.1, 5.5, "搜索与筛选", [
    "关键词 + MySQL FULLTEXT 全文索引",
    "分类下拉筛选 (一级 + 二级)",
    "价格区间 + 商品成色筛选",
    "排序: 价格 / 时间 / 热度",
    "URL 参数同步 + 分页组件(每页20条)",
    "空状态 El-Empty 友好提示"])
sec(s, 0.8, 3.5, 11.5, "商品详情页", [
    "图片画廊(主图+缩略图切换) | 成色标签(全新/几乎全新/轻微痕迹/明显痕迹) | 大号价格展示(+原价删除线)",
    "收藏(Star/StarFilled) | 卖家信息卡片(头像+昵称+信誉分) | 联系卖家(一键跳转聊天) | 加入购物车+立即购买 | HTML安全净化防XSS"])

# === Page 3: Core Flow + Payment ===
s = prs.slides.add_slide(prs.slide_layouts[6])
bar(s); logo(s); pn(s, 3); title(s, "项目演示 — 核心交易流程与支付系统")
t = s.shapes.add_textbox(Inches(0.8), Inches(0.95), Inches(11.5), Inches(0.35))
p = t.text_frame.paragraphs[0]
p.text = "C2C 交易闭环: 注册 → 发布 → 浏览 → 沟通 → 下单 → 支付 → 发货 → 收货 → 评价"
p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = GREEN; p.alignment = PP_ALIGN.CENTER

sec(s, 0.8, 1.5, 5.5, "模拟支付系统 (新增)", [
    "Step 1: 订单摘要 + 选择支付方式",
    "(微信支付 / 支付宝 / 校园卡三种方式)",
    "Step 2: 旋转spinner处理中动画",
    "(1.5~2.5s 模拟支付网关连接)",
    "Step 3: SVG勾选成功动画 + 价格确认",
    "(2秒倒计时自动跳转订单详情)",
    "Step 4: 失败抖动动画 + 重试按钮",
    "购物车结算→自动跳转支付→成功→订单详情"], ACCENT)
sec(s, 7, 1.5, 5.5, "消息与通知系统", [
    "站内信: 双栏聊天(会话列表+气泡面板)",
    "商品页'联系卖家'→ ?to=userId 一键开聊",
    "Enter 快捷发送 + 30s 未读轮询",
    "系统通知: 头部铃铛+红点+Popover",
    "下单/付款/发货/收货自动触发通知",
    "消息发送自动创建对方系统通知"])

# === Page 4: Architecture ===
s = prs.slides.add_slide(prs.slide_layouts[6])
bar(s); logo(s); pn(s, 4); title(s, "核心设计 — 技术架构与安全")
sec(s, 0.8, 1.2, 3.6, "前端 (Vue3 SPA)", [
    "Vue 3 + TypeScript + Vite 5",
    "Element Plus + Pinia + Axios",
    "ECharts 5 数据可视化",
    "CSS 设计令牌 + 全局主题"])
sec(s, 4.8, 1.2, 3.6, "后端 (SpringBoot 3)", [
    "SpringBoot 3.2 + MyBatis-Plus 3.5",
    "Spring Security + JWT 0.12",
    "Knife4j API 文档",
    "RESTful API 30+ 端点"])
sec(s, 8.8, 1.2, 3.8, "数据与部署 (Docker)", [
    "MySQL 8.0 + Redis 7",
    "Docker Compose 4容器编排",
    "Nginx反向代理+Gzip+限流",
    "健康检查+数据卷持久化"])
sec(s, 0.8, 3.3, 5.8, "安全设计", [
    "JWT 双Token: Access 2h + Refresh 7d",
    "Token 黑名单: 登出加入Redis防泄露",
    "三级权限: 公开/认证/ROLE_ADMIN",
    "BCrypt 加密 + 登录频率限制(10次/15min)",
    "验证码类型隔离 + 防暴力破解(5次/15min)"])
sec(s, 7, 3.3, 5.5, "订单状态机 + 数据库设计", [
    "待付款 → 待发货 → 待收货 → 已完成/已取消",
    "悲观锁 selectByIdForUpdate 防双重出售",
    "13张表 · InnoDB · utf8mb4 · 全文索引 · 软删除",
    "N+1优化: 1次JOIN替代N次查询",
    "订单日志 + 商品售出同步 + 分类树Redis缓存"])

# === Page 5-9: Team ===
members = [
    ("张可凡（队长）", "全栈开发 / 项目管理 / 测试", ACCENT,
     ["Docker Compose 编排 MySQL+Redis+SpringBoot+Nginx",
      "Nginx 反向代理 + CORS + Gzip + 限流",
      "容器健康检查 + 数据卷持久化 + 时区修复",
      "管理后台 5 个页面 (Vue3+Element Plus)",
      "Postman 30+ 接口测试 + JUnit 14 单元测试",
      "7 份软件工程文档 + 项目进度管理"]),
    ("阮崇轩", "前端开发", GREEN,
     ["Vue3+TS工程搭建 + CSS设计令牌体系",
      "登录/注册/找回密码 + JWT Token自动刷新",
      "消息系统: 双栏聊天 + 30s轮询 + 通知铃铛",
      "订单详情: 买家卖家双视角 + 支付/发货/收货/评价",
      "UI动效: 骨架屏/过渡动画/回到顶部/滚动导航",
      "后端Bug修复: 验证码隔离/防暴力破解/Token废弃"]),
    ("李硕", "前端开发", GREEN,
     ["首页: Hero横幅+公告轮播+分类导航+商品网格",
      "商品模块: 列表筛选+详情画廊+发布多图上传",
      "购物车+订单列表: Tab切换+买家卖家视角",
      "管理后台 5 个页面 (用户/商品审核/订单/统计/公告)",
      "ECharts 饼图+折线图 + dispose 防内存泄漏",
      "217件种子SQL + 三版图片方案迭代"]),
    ("刘轩霆", "后端开发 A", ACCENT,
     ["Spring Security + JWT 无状态认证架构",
      "商品API: 分页+筛选+全文搜索+批量加载优化",
      "订单核心: 悲观锁+状态机+订单日志+自买自卖校验",
      "后台管理API: 用户/商品/订单/统计 全量接口",
      "MyBatis-Plus配置: 逻辑删除+自动填充+分页插件",
      "Knife4j API文档 + 统一返回封装(Result/PageResult)"]),
    ("胡欣悦", "后端开发 B", ACCENT,
     ["购物车CRUD: 添加+列表+修改数量+去重校验",
      "消息模块: 站内信5端点 + 通知3端点 + 会话聚合",
      "数据统计: 首页统计+分类树Redis缓存+后台趋势",
      "文件上传: MIME校验+大小限制+Docker持久化",
      "Docker运维: Dockerfile多阶段构建+后端部署",
      "MySQL初始化脚本+Redis配置+健康检查"])
]

for i, (name, role, color, items) in enumerate(members):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bar(s); logo(s); pn(s, 5+i)
    t = s.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(11), Inches(0.5))
    p = t.text_frame.paragraphs[0]
    p.text = f"{name} — {role}"; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = TEXT
    sec(s, 0.8, 1.0, 11.5, "核心贡献", items, color)

# === Page 10: Summary ===
s = prs.slides.add_slide(prs.slide_layouts[6])
bar(s); logo(s); pn(s, 10); title(s, "项目成果总结")
sec(s, 0.8, 1.2, 5.5, "技术成果", [
    "前端 18,000+ 行 · 后端 8,000+ 行 · SQL 1,400+ 行",
    "21 个前端页面 · 30+ API 接口 · 13 张数据库表",
    "Docker Compose 一键部署 · 7 份软件工程文档",
    "JWT 双Token + Redis黑名单 + BCrypt加密",
    "完整 C2C 交易闭环 + 消息通知实时化"], GREEN)
sec(s, 7, 1.2, 5.5, "项目亮点", [
    "Token 队列刷新: 多请求并发401只刷新一次",
    "验证码三重防护: 类型隔离+频率限制+防消费",
    "订单悲观锁防双重出售+商品售出状态同步",
    "CSS 设计令牌体系 + Element Plus 全局覆盖",
    "骨架屏加载 + 页面过渡动画 + 滚动感知导航",
    "消息/通知/订单状态三方联动自动推送"], ACCENT)
sec(s, 0.8, 4.2, 11.5, "团队协作", [
    "7 天 · 5 人 · 15+ GitHub Commits · 前后端 0 阻塞交付",
    "2 前端 + 2 后端 + 1 全栈/PM · 每日站会 · 接口文档驱动联调"])

# === Page 11: Thank You ===
s = prs.slides.add_slide(prs.slide_layouts[6])
r = s.shapes.add_shape(1, Inches(0), Inches(0), W, H)
r.fill.solid(); r.fill.fore_color.rgb = DARK; r.line.fill.background()
logo(s)
t = s.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.33), Inches(3))
tf = t.text_frame
p = tf.paragraphs[0]; p.text = "感谢聆听"; p.font.size = Pt(56); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph(); p.text = ""; p.font.size = Pt(16)
p = tf.add_paragraph(); p.text = "转转 — 校园二手物品交易平台"; p.font.size = Pt(22); p.font.color.rgb = RGBColor(0xCC,0xCC,0xCC); p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph(); p.text = ""; p.font.size = Pt(12)
p = tf.add_paragraph(); p.text = "张可凡 · 阮崇轩 · 李硕 · 刘轩霆 · 胡欣悦"; p.font.size = Pt(16); p.font.color.rgb = RGBColor(0xAA,0xAA,0xAA); p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph(); p.text = "华中农业大学 · 2026年7月"; p.font.size = Pt(14); p.font.color.rgb = RGBColor(0x99,0x99,0x99); p.alignment = PP_ALIGN.CENTER

out = "F:/zhuanzhuan/Zhuanzhuan/docs/答辩/转转-软件工程答辩-v2.pptx"
prs.save(out)
print(f"Saved: {out}")
