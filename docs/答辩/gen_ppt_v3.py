"""PPT v3 — Large text, focused content, space for images"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

GREEN = RGBColor(0x2D,0x6A,0x4F)
DARK = RGBColor(0x1B,0x43,0x32)
ACCENT = RGBColor(0xE7,0x6F,0x51)
WHITE = RGBColor(0xFF,0xFF,0xFF)
BG = RGBColor(0xFD,0xFB,0xF7)
TEXT = RGBColor(0x1F,0x29,0x37)
MUTED = RGBColor(0x6B,0x72,0x80)
LIGHT = RGBColor(0xE8,0xF5,0xE9)
WARM = RGBColor(0xFF,0xF8,0xF0)
W = prs.slide_width; H = prs.slide_height
LOGO = "F:/zhuanzhuan/Zhuanzhuan/docs/工作日报/hzau-logo.png"

def L(s): s.shapes.add_picture(LOGO, W-Inches(1.8), Inches(0.15), width=Inches(1.5), height=Inches(0.45))
def B(s):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, Inches(0.04))
    r.fill.solid(); r.fill.fore_color.rgb = GREEN; r.line.fill.background()
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), H-Inches(0.03), W, Inches(0.03))
    r.fill.solid(); r.fill.fore_color.rgb = GREEN; r.line.fill.background()
def PN(s,n):
    t = s.shapes.add_textbox(W-Inches(1.2), H-Inches(0.5), Inches(1), Inches(0.3))
    p = t.text_frame.paragraphs[0]; p.text=f"{n}/20"; p.font.size=Pt(11)
    p.font.color.rgb=RGBColor(0x9C,0xA3,0xAF); p.alignment=PP_ALIGN.RIGHT

def big_title(s,txt):
    t=s.shapes.add_textbox(Inches(0.8),Inches(0.25),Inches(11),Inches(0.7))
    p=t.text_frame.paragraphs[0];p.text=txt;p.font.size=Pt(36);p.font.bold=True;p.font.color.rgb=TEXT

def sub_title(s,txt):
    t=s.shapes.add_textbox(Inches(0.8),Inches(0.9),Inches(11),Inches(0.4))
    p=t.text_frame.paragraphs[0];p.text=txt;p.font.size=Pt(18);p.font.color.rgb=MUTED

def bullet_block(s,x,y,w,items,color=GREEN,size=Pt(20)):
    """Big bullet text — 3-4 items max per block"""
    t=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(len(items)*0.55+0.3))
    tf=t.text_frame; tf.word_wrap=True
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text=item; p.font.size=size; p.font.color.rgb=TEXT; p.space_after=Pt(10)
        p.level=0
    return tf

def section_tag(s,x,y,label,color=GREEN):
    """Colored section label"""
    t=s.shapes.add_textbox(Inches(x),Inches(y),Inches(2),Inches(0.4))
    p=t.text_frame.paragraphs[0];p.text=label;p.font.size=Pt(22);p.font.bold=True;p.font.color.rgb=color

def img_placeholder(s,x,y,w_text,h_text,label="📷 图片区域"):
    """Dashed image placeholder"""
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w_text),Inches(h_text))
    r.fill.solid();r.fill.fore_color.rgb=RGBColor(0xF5,0xF5,0xF5);r.line.color.rgb=RGBColor(0xCC,0xCC,0xCC)
    r.line.width=Pt(1.5);r.line.dash_style=2
    t=s.shapes.add_textbox(Inches(x),Inches(y+h_text/2-0.3),Inches(w_text),Inches(0.6))
    p=t.text_frame.paragraphs[0];p.text=label;p.font.size=Pt(16)
    p.font.color.rgb=RGBColor(0xBB,0xBB,0xBB);p.alignment=PP_ALIGN.CENTER

def stat_card(s,x,y,num,label,color=GREEN):
    r=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(2.6),Inches(1.2))
    r.fill.solid();r.fill.fore_color.rgb=LIGHT;r.line.fill.background();r.adjustments[0]=0.1
    t=s.shapes.add_textbox(Inches(x),Inches(y+0.1),Inches(2.6),Inches(0.6))
    p=t.text_frame.paragraphs[0];p.text=num;p.font.size=Pt(36);p.font.bold=True
    p.font.color.rgb=color;p.alignment=PP_ALIGN.CENTER
    t=s.shapes.add_textbox(Inches(x),Inches(y+0.7),Inches(2.6),Inches(0.4))
    p=t.text_frame.paragraphs[0];p.text=label;p.font.size=Pt(14);p.font.color.rgb=MUTED;p.alignment=PP_ALIGN.CENTER

# ==============================
# Page 1: COVER
# ==============================
s=prs.slides.add_slide(prs.slide_layouts[6])
r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0),Inches(0),W,H)
r.fill.solid();r.fill.fore_color.rgb=DARK;r.line.fill.background()
L(s)
t=s.shapes.add_textbox(Inches(1),Inches(1.5),Inches(11.33),Inches(4))
tf=t.text_frame
p=tf.paragraphs[0];p.text="转转";p.font.size=Pt(80);p.font.bold=True;p.font.color.rgb=WHITE;p.alignment=PP_ALIGN.CENTER
p=tf.add_paragraph();p.text="校园二手物品交易平台";p.font.size=Pt(30);p.font.color.rgb=RGBColor(0xCC,0xCC,0xCC);p.alignment=PP_ALIGN.CENTER
p=tf.add_paragraph();p.text="";p.font.size=Pt(20)
p=tf.add_paragraph();p.text="软件工程实训 · 课程设计答辩";p.font.size=Pt(24);p.font.color.rgb=RGBColor(0xAA,0xAA,0xAA);p.alignment=PP_ALIGN.CENTER
p=tf.add_paragraph();p.text="";p.font.size=Pt(16)
p=tf.add_paragraph();p.text="张可凡 · 阮崇轩 · 李硕 · 刘轩霆 · 胡欣悦";p.font.size=Pt(18);p.font.color.rgb=RGBColor(0x88,0x88,0x88);p.alignment=PP_ALIGN.CENTER
p=tf.add_paragraph();p.text="华中农业大学 · 2026年7月";p.font.size=Pt(16);p.font.color.rgb=RGBColor(0x77,0x77,0x77);p.alignment=PP_ALIGN.CENTER

# ==============================
# Page 2: OVERVIEW + STATS
# ==============================
s=prs.slides.add_slide(prs.slide_layouts[6]);B(s);L(s);PN(s,2)
big_title(s,"项目概述")
sub_title(s,"校园 C2C 二手交易平台 — 前后端分离 · Docker 容器化 · 完整交易闭环")
stat_card(s,0.8,1.7,"26,000+","总代码行数")
stat_card(s,3.6,1.7,"30+","REST API")
stat_card(s,6.4,1.7,"21","前端页面")
stat_card(s,9.2,1.7,"217","商品数据")
section_tag(s,0.8,3.3,"核心技术栈")
bullet_block(s,0.8,3.8,5.5,[
    "Vue 3 + TypeScript + Vite 5 + Element Plus",
    "SpringBoot 3.2 + MyBatis-Plus + Security + JWT",
    "MySQL 8.0（全文索引） + Redis 7（缓存/黑名单）",
    "Docker Compose 五容器编排 + Nginx 反向代理"
])
img_placeholder(s,7,3.8,5.5,2.5,"📷 架构图或截图")

# ==============================
# Page 3: CORE FEATURES (split into 2 pages)
# ==============================
s=prs.slides.add_slide(prs.slide_layouts[6]);B(s);L(s);PN(s,3)
big_title(s,"浏览与交易")
section_tag(s,0.8,1.3,"商品浏览")
bullet_block(s,0.8,1.75,5.5,[
    "首页 Hero 横幅 + 实时统计（17用户/217商品）",
    "分类 Tab 切换 · 全文搜索 · 价格/成色筛选",
    "商品详情：图片画廊 + 卖家信息卡片",
    "联系卖家按钮 → 一键跳转聊天窗口"
])
img_placeholder(s,7,1.3,5.5,2.3,"📷 首页截图")
section_tag(s,0.8,4.0,"交易流程")
bullet_block(s,0.8,4.45,5.5,[
    "下单时弹出地址选择（默认地址高亮）",
    "模拟支付：3 种方式 + spinner + 成功动画",
    "发货：物流公司卡片选择 + 单号录入"
])
img_placeholder(s,7,4.0,5.5,2.3,"📷 支付页面截图")

# ==============================
# Page 4: MESSAGES & ORDERS
# ==============================
s=prs.slides.add_slide(prs.slide_layouts[6]);B(s);L(s);PN(s,4)
big_title(s,"消息与订单")
section_tag(s,0.8,1.3,"实时消息")
bullet_block(s,0.8,1.75,5.5,[
    "双栏聊天：会话列表 + 聊天气泡",
    "每 5 秒自动拉取新消息（无需刷新）",
    "头部铃铛通知 · 点击直接跳转对应页面",
    "发消息/下单/付款/发货 自动推送通知"
])
img_placeholder(s,7,1.3,5.5,2.3,"📷 聊天截图")
section_tag(s,0.8,4.0,"订单管理")
bullet_block(s,0.8,4.45,5.5,[
    "订单状态进度条（待付款→待发货→待收货→完成）",
    "买家：支付 · 确认收货 · 星级评价",
    "卖家：录入物流（8家快递） · 查看待发货",
])
img_placeholder(s,7,4.0,5.5,2.3,"📷 订单截图")

# ==============================
# Page 5: ADMIN
# ==============================
s=prs.slides.add_slide(prs.slide_layouts[6]);B(s);L(s);PN(s,5)
big_title(s,"管理后台")
section_tag(s,0.8,1.3,"五大管理模块")
bullet_block(s,0.8,1.75,5.5,[
    "用户管理：列表 · 启用/禁用 · 角色分配",
    "商品审核：通过/驳回 · 违规下架 · 编辑",
    "订单管理：状态筛选 · 标记操作 · 详情",
    "数据统计：ECharts 饼图 + 折线图 · 自适应"
])
bullet_block(s,7,1.75,5.5,[
    "公告管理：发布/草稿 · 编辑 · 删除",
    "侧边栏导航 · 路由鉴权 · 响应式折叠",
    "非管理员自动重定向到首页"
])
img_placeholder(s,0.8,3.5,11.5,3.2,"📷 管理后台截图")

# ==============================
# Page 6: SECURITY
# ==============================
s=prs.slides.add_slide(prs.slide_layouts[6]);B(s);L(s);PN(s,6)
big_title(s,"安全架构")
section_tag(s,0.8,1.3,"认证与授权",ACCENT)
bullet_block(s,0.8,1.75,5.5,[
    "JWT 双 Token：Access 2h + Refresh 7d",
    "请求队列机制：多请求 401 只刷新一次",
    "Token 黑名单：登出即失效（Redis）",
    "三级权限：公开 · 认证 · 管理员"
],ACCENT)
section_tag(s,7,1.3,"防护措施",ACCENT)
bullet_block(s,7,1.75,5.5,[
    "验证码类型隔离（注册/重置不可互用）",
    "频率限制：登录/验证码/密码重置三层",
    "BCrypt 密码哈希 + XSS 过滤 + MIME 校验",
    "软删除 + 敏感信息脱敏"
],ACCENT)
section_tag(s,0.8,3.8,"订单安全",GREEN)
bullet_block(s,0.8,4.25,5.5,[
    "悲观锁 selectForUpdate 防双重出售",
    "自买自卖校验 · 订单日志全追溯",
    "13 张表 · InnoDB · utf8mb4 · 全文索引"
])
img_placeholder(s,7,3.8,5.5,2.2,"📷 安全架构图")

# ==============================
# Page 7: DEPLOYMENT
# ==============================
s=prs.slides.add_slide(prs.slide_layouts[6]);B(s);L(s);PN(s,7)
big_title(s,"Docker 部署架构")
section_tag(s,0.8,1.3,"五容器编排")
bullet_block(s,0.8,1.75,5.5,[
    "Nginx（Vue3 SPA → :3000）",
    "SpringBoot 3.2（Java 17 → :8080）",
    "MySQL 8.0 · Redis 7 · phpMyAdmin（:3001）",
    "Docker Network Bridge · 内部通信"
])
sub_title(s,"一键启动：docker-compose up -d --build")
section_tag(s,7,1.3,"运维特性")
bullet_block(s,7,1.75,5.5,[
    "健康检查依赖启动顺序",
    "4 个 MySQL 自动初始化脚本",
    "数据卷持久化（mysql/redis/upload）",
    "TZ=Asia/Shanghai 时区修复"
])
img_placeholder(s,0.8,3.6,11.5,3.2,"📷 Docker 架构图或容器运行截图")

# ==============================
# Pages 6-10: TEAM MEMBERS (5 members)
# ==============================
members = [
    ("全栈开发 / 项目管理 / 测试", ACCENT, [
        "编写 docker-compose.yml 五容器编排 + Nginx 配置",
        "管理后台前端 5 个页面（用户/商品/订单/统计/公告）",
        "Postman 30+ 接口测试 + 前后端全量联调",
        "7 份软件工程文档 + README + 项目进度管理",
        "发现并修复中文编码 · 时区偏移 · JWT 刷新冲突等关键 Bug"
    ]),
    ("前端开发", GREEN, [
        "Vue3 工程搭建 + CSS 设计令牌体系 + Element Plus 全局覆盖",
        "登录 / 注册 / 找回密码 · Axios Token 自动刷新队列",
        "消息系统：双栏聊天 + 5s 实时轮询 + 通知铃铛 Popover",
        "订单详情：买家卖家双视角 · 物流录入 · 评价打分",
        "UI 动效：骨架屏 · 过渡动画 · 回到顶部 · 滚动导航"
    ]),
    ("前端开发", GREEN, [
        "首页：Hero 横幅 + 公告轮播 + 分类导航 + 商品网格",
        "商品模块：列表筛选 · 详情画廊 · 发布多图上传",
        "购物车 + 订单列表（Tab 切换 + 双视角）",
        "管理后台 5 页面 + ECharts 饼图/折线图可视化",
        "217 件种子数据 SQL + 三版图片方案迭代"
    ]),
    ("后端开发 A", ACCENT, [
        "Spring Security + JWT 认证 · 三级权限 · 频率限制",
        "商品 API：分页/搜索/筛选/排序 · N+1 批量优化",
        "订单核心：悲观锁 · 状态机 · 订单日志 · 自买自卖校验",
        "13 张数据库表设计（ENUM 约束 · 全文索引 · 软删除）",
        "后台管理 API · Knife4j 文档 · 14 单元测试全通过"
    ]),
    ("后端开发 B", ACCENT, [
        "购物车 CRUD + 去重 · 消息/通知模块（8 端点）",
        "数据统计（首页 · 分类树 Redis 缓存 · 后台趋势）",
        "文件上传（MIME 校验 · Docker Volume 持久化）",
        "Docker 运维：后端多阶段构建 · MySQL 初始化脚本",
        "订单/消息自动通知触发 + phpMyAdmin 集成"
    ])
]

for i,(role,color,items) in enumerate(members):
    s=prs.slides.add_slide(prs.slide_layouts[6]);B(s);L(s);PN(s,8+i)
    big_title(s,role)
    # Left: key contributions (big text)
    bullet_block(s,0.8,1.3,7.5,items,color,Pt(20))
    # Right: image placeholder
    img_placeholder(s,8.8,1.3,3.8,5.0,f"📷 截图")

# ==============================
# Page 13: FRONTEND DESIGN
# ==============================
s=prs.slides.add_slide(prs.slide_layouts[6]);B(s);L(s);PN(s,13)
big_title(s,"前端设计体系")
section_tag(s,0.8,1.3,"CSS 设计令牌")
bullet_block(s,0.8,1.75,5.5,[
    "主色 #2D6A4F（森林绿）+ 强调 #E76F51（暖橙）",
    "8 个圆角/阴影变量 + 品牌色全局统一",
    "Element Plus 全局覆盖 8 类组件",
    "21 个页面风格 100% 一致"
])
section_tag(s,7,1.3,"动效系统")
bullet_block(s,7,1.75,5.5,[
    "骨架屏 Shimmer · 路由过渡 fade/slide",
    "滚动感知毛玻璃导航 · 按钮微交互",
    "卡片 hover 放大 · fadeInUp 入场动画",
    "通知 Popover · 支付勾选/抖动动画"
])
section_tag(s,0.8,3.6,"用户体验")
bullet_block(s,0.8,4.05,11.5,[
    "实时聊天 5s 轮询 · 通知 30s 轮询 · 地址选择对话框 · 物流单号可复制",
    "空状态友好引导 · 错误分类提示 · 403 过期自动清除 · 订单状态可视化进度条"
])

# ==============================
# Page 14: DATABASE
# ==============================
s=prs.slides.add_slide(prs.slide_layouts[6]);B(s);L(s);PN(s,14)
big_title(s,"数据库设计")
section_tag(s,0.8,1.3,"13 张核心表")
bullet_block(s,0.8,1.75,5.5,[
    "user · product · category · product_image",
    "order · order_log · cart · address",
    "message · notification · favorite · review · announcement",
    "ENUM 约束 · 全文索引 · 软删除 · utf8mb4"
])
section_tag(s,7,1.3,"订单状态机",ACCENT)
bullet_block(s,7,1.75,5.5,[
    "待付款 → 待发货 → 待收货 → 已完成 / 已取消",
    "悲观锁 selectForUpdate 防双重出售",
    "order_log 记录每次状态变更操作人",
],ACCENT)
img_placeholder(s,0.8,3.5,11.5,3.2,"📷 ER 图或数据库截图")

# ==============================
# Page 15: ARCHITECTURE DETAIL
# ==============================
s=prs.slides.add_slide(prs.slide_layouts[6]);B(s);L(s);PN(s,15)
big_title(s,"技术架构详解")
section_tag(s,0.8,1.3,"请求链路")
bullet_block(s,0.8,1.75,11.5,[
    "浏览器 :3000 → Nginx → /api/* → SpringBoot :8080 → MySQL/Redis → JSON 响应",
    "Vite 开发代理 → 生产 Nginx 反向代理 → Gzip 压缩 → 限流 30r/s → CORS 白名单"
])
section_tag(s,0.8,2.8,"关键设计")
bullet_block(s,0.8,3.2,5.5,[
    "前后端分离：独立开发/部署，JSON API 通信",
    "JWT 无状态：Access 2h + Refresh 7d",
    "统一封装：Result / PageResult JSON 标准化",
    "全局异常：@RestControllerAdvice 统一处理"
])
bullet_block(s,7,3.2,5.5,[
    "MyBatis-Plus：逻辑删除 · 自动填充 · 分页",
    "Redis：验证码 · Token黑名单 · 分类树缓存",
    "Knife4j：自动生成 Swagger API 文档",
    "14 单元测试 · 30+ Postman 接口测试"
])
img_placeholder(s,0.8,5.3,11.5,1.5,"📷 架构图")

# ==============================
# Page 16: TESTING
# ==============================
s=prs.slides.add_slide(prs.slide_layouts[6]);B(s);L(s);PN(s,16)
big_title(s,"测试与质量保证")
section_tag(s,0.8,1.3,"测试覆盖")
bullet_block(s,0.8,1.75,5.5,[
    "JUnit 后端单元测试：14 个 · 全部通过",
    "Postman 接口测试：30+ 用例 · 7 大模块",
    "Vitest 前端测试：8 个 · 全部通过",
    "前后端全量联调：逐接口验证"
])
section_tag(s,7,1.3,"Bug 修复记录",ACCENT)
bullet_block(s,7,1.75,5.5,[
    "验证码类型隔离 + 防消费 + 防暴力破解",
    "密码重置 Token 提前删除 Bug",
    "Docker 时区偏移（UTC→Asia/Shanghai）",
    "旧 Token 全站 403 · 卖家浏览量误计"
],ACCENT)
img_placeholder(s,0.8,3.5,11.5,3.2,"📷 测试报告或覆盖率截图")

# ==============================
# Page 17: UI/UX HIGHLIGHTS
# ==============================
s=prs.slides.add_slide(prs.slide_layouts[6]);B(s);L(s);PN(s,17)
big_title(s,"UI/UX 设计亮点")
bullet_block(s,0.8,1.3,5.5,[
    "CSS 设计令牌：15 个变量全局统一品牌色",
    "Element Plus 全组件覆盖：8 类组件统一风格",
    "骨架屏 Shimmer：首页 8 卡片先展示",
    "路由过渡：fade/slide mode=out-in"
])
bullet_block(s,7,1.3,5.5,[
    "按钮微交互：hover 上浮 1px + 阴影",
    "滚动导航：backdrop-filter blur(12px)",
    "回到顶部按钮：400px 弹性出现",
    "卡片入场：fadeInUp + animationDelay",
])
img_placeholder(s,0.8,3.0,5.5,2.5,"📷 动效录屏")
img_placeholder(s,7,3.0,5.5,2.5,"📷 组件库截图")

# ==============================
# Page 18: DATA FLOW
# ==============================
s=prs.slides.add_slide(prs.slide_layouts[6]);B(s);L(s);PN(s,18)
big_title(s,"数据与通知流程")
section_tag(s,0.8,1.3,"消息通知链路")
bullet_block(s,0.8,1.75,11.5,[
    "用户发消息 → MessageServiceImpl → 写入 message 表 + 创建 Notification → 对方轮询获取",
    "创建订单 → 通知卖家「有新订单」 → 买家付款 → 通知卖家「请发货」 → 卖家发货 → 通知买家「已发货」",
])
section_tag(s,0.8,2.8,"数据流")
bullet_block(s,0.8,3.2,5.5,[
    "分类树：MySQL 递归查询 → Redis 缓存",
    "验证码：生成 → Redis(5min) → 校验 → 删除",
    "Token 黑名单：登出 → Redis → 过滤器检查",
])
bullet_block(s,7,3.2,5.5,[
    "文件上传：MIME 校验 → 本地存储 → Docker Volume",
    "浏览量：非卖家浏览 → UPDATE view_count+1",
    "定时轮询：消息 5s · 通知 30s · 自动刷新",
])
img_placeholder(s,0.8,4.8,11.5,1.8,"📷 数据流图或通知截图")

# ==============================
# Page 19: RESULTS
# ==============================
s=prs.slides.add_slide(prs.slide_layouts[6]);B(s);L(s);PN(s,19)
big_title(s,"项目成果总结")
section_tag(s,0.8,1.3,"交付成果")
bullet_block(s,0.8,1.75,5.5,[
    "前端 18,000+ 行 · 后端 8,000+ 行 · SQL 1,400+ 行",
    "30+ API · 21 页面 · 13 表 · 217 商品数据",
    "Docker 五容器一键部署 + 自动初始化",
    "7 份软件工程文档 · 20 页答辩 PPT"
])
section_tag(s,7,1.3,"项目亮点",ACCENT)
bullet_block(s,0.8,3.3,5.5,[
    "完整 C2C 闭环（注册→交易→评价）",
    "JWT 双 Token + 三层安全防护",
    "消息通知实时化 + 订单自动通知",
    "Docker 一键部署 + 自动初始化"
])
section_tag(s,7,2.9,"项目亮点",ACCENT)
bullet_block(s,7,3.3,5.5,[
    "Token 队列刷新：多请求 401 只刷一次",
    "验证码三重防护：隔离 + 限频 + 防消费",
    "N+1 查询优化 + 悲观锁防双重出售",
    "CSS 令牌体系 + 骨架屏 + 动效系统"
],ACCENT)
img_placeholder(s,0.8,3.8,11.5,2.5,"📷 项目全景截图或架构图")

# ==============================
# Page 20: THANKS
# ==============================
s=prs.slides.add_slide(prs.slide_layouts[6]);PN(s,20)
r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0),Inches(0),W,H)
r.fill.solid();r.fill.fore_color.rgb=DARK;r.line.fill.background()
L(s)
t=s.shapes.add_textbox(Inches(1),Inches(2),Inches(11.33),Inches(3))
tf=t.text_frame
p=tf.paragraphs[0];p.text="感谢聆听";p.font.size=Pt(64);p.font.bold=True;p.font.color.rgb=WHITE;p.alignment=PP_ALIGN.CENTER
p=tf.add_paragraph();p.text="";p.font.size=Pt(20)
p=tf.add_paragraph();p.text="转转 — 校园二手物品交易平台";p.font.size=Pt(26);p.font.color.rgb=RGBColor(0xCC,0xCC,0xCC);p.alignment=PP_ALIGN.CENTER
p=tf.add_paragraph();p.text="";p.font.size=Pt(14)
p=tf.add_paragraph();p.text="华中农业大学 · 2026年7月";p.font.size=Pt(18);p.font.color.rgb=RGBColor(0x99,0x99,0x99);p.alignment=PP_ALIGN.CENTER

out="F:/zhuanzhuan/Zhuanzhuan/docs/答辩/转转-答辩-v3.pptx"
prs.save(out)
print(f"Saved: {out} ({prs.slides.__len__()} slides)")
